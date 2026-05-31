"""
build_pptx.py
============
Genera la presentación de defensa de la tesis (PowerPoint) con estética moderna
(fondo blanco, tipografía limpia, azul acento) e incrustando las figuras reales.

Salida: outputs/Defensa_Tesis_Francisco_Rietta.pptx
"""
from pathlib import Path
import json
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

ROOT = Path(__file__).resolve().parent
FIG = ROOT / "outputs" / "figures"
PLAT = ROOT.parent / "tesis-plataforma"
OUT = ROOT / "outputs" / "Defensa_Tesis_Francisco_Rietta.pptx"

INK = RGBColor(0x1D, 0x1D, 0x1F)
INK2 = RGBColor(0x6E, 0x6E, 0x73)
BLUE = RGBColor(0x00, 0x71, 0xE3)
COPPER = RGBColor(0xC9, 0x71, 0x2D)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT = RGBColor(0xFB, 0xFB, 0xFD)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def _data():
    try:
        return json.loads((PLAT / "data.json").read_text(encoding="utf-8"))
    except Exception:
        return {}


D = _data()


def bg(slide, color=WHITE):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def textbox(slide, x, y, w, h, text, size, color=INK, bold=False, align=PP_ALIGN.LEFT,
            font="Calibri Light", anchor=MSO_ANCHOR.TOP, italic=False, line_spacing=1.05):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    lines = text.split("\n")
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.line_spacing = line_spacing
        r = p.add_run(); r.text = ln
        r.font.size = Pt(size); r.font.bold = bold; r.font.italic = italic
        r.font.color.rgb = color; r.font.name = font
    return tb


def bullets(slide, x, y, w, h, items, size=20, color=INK, gap=10):
    tb = slide.shapes.add_textbox(x, y, w, h); tf = tb.text_frame; tf.word_wrap = True
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap); p.line_spacing = 1.1
        bold_part = None
        if isinstance(it, tuple):
            bold_part, it = it
        r0 = p.add_run(); r0.text = "•  "; r0.font.size = Pt(size); r0.font.color.rgb = BLUE
        if bold_part:
            rb = p.add_run(); rb.text = bold_part + " "
            rb.font.size = Pt(size); rb.font.bold = True; rb.font.color.rgb = color; rb.font.name = "Calibri"
        r = p.add_run(); r.text = it
        r.font.size = Pt(size); r.font.color.rgb = color; r.font.name = "Calibri"
    return tb


def accent_bar(slide, x=Inches(0.7), y=Inches(1.5), w=Inches(0.9)):
    from pptx.enum.shapes import MSO_SHAPE
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, Pt(5))
    sh.fill.solid(); sh.fill.fore_color.rgb = BLUE; sh.line.fill.background()


def kicker_title(slide, kicker, title, tsize=34):
    textbox(slide, Inches(0.7), Inches(0.55), Inches(11.9), Inches(0.5),
            kicker.upper(), 14, BLUE, bold=True, font="Calibri")
    accent_bar(slide, y=Inches(1.05), w=Inches(0.7))
    textbox(slide, Inches(0.7), Inches(1.15), Inches(11.9), Inches(1.2),
            title, tsize, INK, bold=True)


def img(slide, path, x, y, w):
    p = FIG / path
    if p.exists():
        slide.shapes.add_picture(str(p), x, y, width=w)


def footer(slide, n):
    textbox(slide, Inches(0.7), Inches(7.0), Inches(8), Inches(0.4),
            "Impacto macroeconómico en el sector cobre · Universidad Señor de Sipán",
            10, INK2, font="Calibri")
    textbox(slide, Inches(12.2), Inches(7.0), Inches(0.8), Inches(0.4),
            str(n), 10, INK2, align=PP_ALIGN.RIGHT, font="Calibri")


# ---------------- Slides ----------------
def s_title():
    s = prs.slides.add_slide(BLANK); bg(s)
    logo = next((c for c in (ROOT/"assets"/"logo_uss.png", PLAT/"logo_uss.png") if c.exists()), None)
    if logo:
        s.shapes.add_picture(str(logo), Inches(5.9), Inches(0.6), height=Inches(1.2))
    textbox(s, Inches(1), Inches(2.0), Inches(11.3), Inches(0.5),
            "UNIVERSIDAD SAN SEBASTIÁN · MAGÍSTER EN DATA SCIENCE", 14, BLUE, bold=True,
            align=PP_ALIGN.CENTER, font="Calibri")
    textbox(s, Inches(1), Inches(2.7), Inches(11.3), Inches(2.2),
            "Impacto de las variables macroeconómicas en los retornos accionarios de las empresas "
            "de cobre con exposición a Chile", 34, INK, bold=True, align=PP_ALIGN.CENTER)
    textbox(s, Inches(1.5), Inches(4.9), Inches(10.3), Inches(0.7),
            "Una comparación entre el mercado bursátil internacional y el chileno · 2004–2024",
            20, INK2, align=PP_ALIGN.CENTER, italic=True)
    textbox(s, Inches(1), Inches(6.2), Inches(11.3), Inches(0.6),
            "Francisco Rietta   ·   Profesor guía: [Nombre]   ·   Santiago, 2026",
            15, INK, align=PP_ALIGN.CENTER, font="Calibri")


def s_problema():
    s = prs.slides.add_slide(BLANK); bg(s)
    kicker_title(s, "Contexto y problema", "Por qué importa")
    bullets(s, Inches(0.7), Inches(2.3), Inches(11.9), Inches(4.2), [
        ("El cobre", "es la columna de la economía chilena: define exportaciones, fisco y tipo de cambio."),
        ("La pregunta abierta:", "¿cuánto, cómo y dónde se transmite el shock macro a las acciones de cobre?"),
        ("Vacío:", "la literatura macro→retornos se concentra en mercados desarrollados e índices agregados; "
         "el cobre, a nivel de empresas con exposición a Chile, está sin estudiar."),
        ("Nadie", "ha comparado formalmente cómo el mercado internacional y el chileno precian el mismo cobre."),
    ], size=21)
    footer(s, 2)


def s_objetivo():
    s = prs.slides.add_slide(BLANK); bg(s)
    kicker_title(s, "Objetivo y preguntas", "Qué se busca responder")
    textbox(s, Inches(0.7), Inches(2.2), Inches(11.9), Inches(1.4),
            "Cuantificar y comparar el impacto de las variables macroeconómicas globales y "
            "nacionales sobre los retornos accionarios de las empresas de cobre con exposición a "
            "Chile (2004–2024), contrastando su transmisión entre el mercado internacional y el "
            "chileno.", 22, INK, italic=True)
    bullets(s, Inches(0.7), Inches(4.0), Inches(11.9), Inches(2.6), [
        ("H1:", "el cobre es el principal determinante (efecto positivo)."),
        ("H6:", "los factores globales dominan a los locales."),
        ("H7:", "el mercado internacional incorpora el shock más completamente que el chileno."),
    ], size=21)
    footer(s, 3)


def s_diseno():
    s = prs.slides.add_slide(BLANK); bg(s)
    kicker_title(s, "Diseño", "Triangulación por tres muestras")
    from pptx.enum.shapes import MSO_SHAPE
    cards = [("B · mercado internacional", "Cobre puro: Antofagasta, BHP, Anglo, Lundin, Teck.\nPanel de efectos fijos.", BLUE),
             ("A · mercado chileno", "Cobre puro local: Pucobre (Bolsa de Santiago).\nSerie de tiempo.", COPPER),
             ("C · sector minero (vehículo)", "Pucobre, CAP, SQM, Molymet.\nPanel; controla por cada commodity.", RGBColor(0xFF,0x9F,0x0A))]
    x = Inches(0.7)
    for t, d, col in cards:
        card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(2.4), Inches(3.85), Inches(3.2))
        card.fill.solid(); card.fill.fore_color.rgb = LIGHT; card.line.color.rgb = RGBColor(0xE0,0xE0,0xE5)
        card.shadow.inherit = False
        textbox(s, x+Inches(0.25), Inches(2.65), Inches(3.4), Inches(0.7), t, 17, col, bold=True)
        textbox(s, x+Inches(0.25), Inches(3.6), Inches(3.4), Inches(1.8), d, 15, INK2)
        x += Inches(4.05)
    textbox(s, Inches(0.7), Inches(5.9), Inches(11.9), Inches(0.8),
            "El foco es el cobre. La idea nueva: ¿el mismo subyacente se precia distinto según el mercado?",
            16, INK, italic=True)
    footer(s, 4)


def s_metodo():
    s = prs.slides.add_slide(BLANK); bg(s)
    kicker_title(s, "Metodología", "Una secuencia econométrica")
    bullets(s, Inches(0.7), Inches(2.2), Inches(11.9), Inches(4.4), [
        ("OE1 ·", "Raíz unitaria (ADF/PP/KPSS/Zivot-Andrews) → orden de integración."),
        ("OE2 ·", "Panel de efectos fijos con errores Driscoll-Kraay (dependencia de sección cruzada)."),
        ("OE3 ·", "Cointegración: ARDL, Johansen y Gregory-Hansen (con quiebre)."),
        ("OE4 ·", "VAR/VECM, impulso-respuesta y descomposición de varianza (FEVD)."),
        ("OE5 ·", "Fechado del ciclo (Bry-Boschan) e interacciones por fase."),
        ("OE6 ·", "Comparación entre mercados sobre las tres muestras."),
    ], size=20, gap=12)
    footer(s, 5)


def s_fig(kicker, title, fig, caption, n, note=None):
    s = prs.slides.add_slide(BLANK); bg(s)
    kicker_title(s, kicker, title, tsize=30)
    img(s, fig, Inches(0.9), Inches(2.1), Inches(8.2))
    textbox(s, Inches(9.4), Inches(2.4), Inches(3.4), Inches(4),
            caption, 18, INK)
    if note:
        textbox(s, Inches(9.4), Inches(5.6), Inches(3.4), Inches(1.2), note, 14, COPPER, bold=True)
    footer(s, n)


def s_resultados_clave():
    s = prs.slides.add_slide(BLANK); bg(s)
    kicker_title(s, "Resultados", "Hallazgos centrales")
    bullets(s, Inches(0.7), Inches(2.2), Inches(11.9), Inches(4.4), [
        ("Cobre +0,57***", "y robusto a corrección por pruebas múltiples (FDR): es el driver (H1)."),
        ("Dominancia global (H6):", "cobre + VIX explican ~32% de la varianza vs ~5% de los locales."),
        ("Cointegración con quiebre (2008):", "la relación de largo plazo existe pero se reconfiguró con la crisis (Gregory-Hansen)."),
        ("Volatilidad asimétrica:", "GJR-GARCH detecta efecto apalancamiento (γ=0,22)."),
        ("Comparación de mercados (H7):", "el internacional incorpora algo más completamente el shock global."),
    ], size=20, gap=11)
    footer(s, 9)


def s_conclusion():
    s = prs.slides.add_slide(BLANK); bg(s)
    bg(s, INK)
    textbox(s, Inches(1), Inches(2.4), Inches(11.3), Inches(0.5),
            "EL HALLAZGO MEMORABLE", 14, BLUE, bold=True, align=PP_ALIGN.CENTER, font="Calibri")
    textbox(s, Inches(1.2), Inches(3.1), Inches(10.9), Inches(2.4),
            "El mismo cobre se precia de forma comparable en distintos mercados, pero el "
            "internacional incorpora más completamente los shocks globales — y la relación de "
            "largo plazo con el cobre se quebró con la crisis de 2008.",
            26, WHITE, bold=True, align=PP_ALIGN.CENTER)
    textbox(s, Inches(1), Inches(6.4), Inches(11.3), Inches(0.5),
            "Gracias.", 20, RGBColor(0xC0,0xC0,0xC8), align=PP_ALIGN.CENTER, italic=True)


def s_vacio():
    s = prs.slides.add_slide(BLANK); bg(s)
    kicker_title(s, "Contribución", "El vacío que llena la tesis")
    bullets(s, Inches(0.7), Inches(2.2), Inches(11.9), Inches(4.4), [
        ("Cobre ↔ peso", "está bien documentado (Chen-Rogoff 2003; Pincheira-Hardy 2019): el peso es una commodity currency."),
        ("Cobre ↔ bolsa chilena agregada", "lo estudió Zurita-Fuentes-Gregoire (2005), pero a nivel de índice y hasta 2003."),
        ("Cobre ↔ acciones mineras", "solo se ha modelado para NY, Toronto y Australia — excluyendo a Chile."),
        ("El vacío:", "nadie estima el efecto a nivel de empresas de cobre con exposición a Chile, ni compara el mercado internacional con el chileno para 2004–2024."),
    ], size=20, gap=12)
    footer(s, 4)


def s_robustez():
    s = prs.slides.add_slide(BLANK); bg(s)
    kicker_title(s, "Robustez", "La evidencia resiste el escrutinio")
    bullets(s, Inches(0.7), Inches(2.2), Inches(11.9), Inches(4.4), [
        ("CD-Pesaran = 24,5:", "hay dependencia de sección cruzada → se usan errores Driscoll-Kraay."),
        ("CIPS (Pesaran 2007):", "raíz unitaria en panel robusta a esa dependencia confirma precios I(1), retornos I(0)."),
        ("Corrección FDR:", "cobre, VIX y tipo de cambio siguen significativos — no son falsos positivos."),
        ("GJR-GARCH:", "efecto apalancamiento (γ=0,22); Local Projections cross-validan la IRF."),
        ("Subperíodos:", "el efecto del cobre se mantiene (0,56 en 2004-19 → 0,75 en 2020-24)."),
    ], size=19, gap=11)
    footer(s, 12)


def s_predictor():
    s = prs.slides.add_slide(BLANK); bg(s)
    kicker_title(s, "Extensión predictiva", "Explicar no es predecir")
    bullets(s, Inches(0.7), Inches(2.3), Inches(11.9), Inches(3.4), [
        ("Los factores macro EXPLICAN", "el retorno contemporáneo del cobre (R² within 0,24, cobre ***)."),
        ("Pero NO lo ANTICIPAN:", "el R² fuera de muestra es ≈0 o negativo en todos los modelos a un mes."),
        ("Lectura:", "coherente con la eficiencia de mercado (forma débil); justifica el enfoque explicativo."),
    ], size=21, gap=13)
    textbox(s, Inches(0.7), Inches(6.0), Inches(11.9), Inches(0.8),
            "Una versión interactiva del predictor está disponible en la plataforma web del proyecto.",
            16, INK2, italic=True)
    footer(s, 13)


def build():
    s_title(); s_problema(); s_objetivo(); s_vacio(); s_diseno(); s_metodo()
    s_fig("Resultados · OE5", "El ciclo del cobre", "fig_ciclo_cobre.png",
          "Fases de expansión/contracción fechadas con Bry-Boschan, base del análisis por régimen.", 6)
    s_fig("Resultados · OE2", "Sensibilidad por factor", "fig_coeficientes.png",
          "Panel de efectos fijos (Driscoll-Kraay). El cobre domina; el tipo de cambio pesa en negativo.", 7,
          note="Cobre +0,57***")
    s_fig("Resultados · OE4", "Descomposición de varianza", "fig_fevd.png",
          "Los shocks globales (cobre + VIX ≈ 32%) superan ampliamente a los locales (≈ 5%).", 8,
          note="Evidencia de H6")
    s_resultados_clave()
    s_fig("Resultados · OE6", "Comparación entre mercados", "fig_triangulacion.png",
          "β del cobre y dominancia global por muestra. Cobre puro: B≈A; menor en la minería mixta (C).", 10)
    s_fig("Resultados · OE4", "Impulso-respuesta y proyecciones locales", "fig_irf.png",
          "La respuesta positiva e inmediata al shock del cobre se confirma con Local Projections (Jordà).", 11)
    s_robustez()
    s_predictor()
    s_conclusion()
    OUT.parent.mkdir(exist_ok=True)
    prs.save(OUT)
    print(f"[ok] {OUT}  ({OUT.stat().st_size//1024} KB, {len(prs.slides.__iter__.__self__._sldIdLst)} slides)")


if __name__ == "__main__":
    build()
